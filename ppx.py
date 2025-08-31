from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import requests
from io import BytesIO

# Vine green as specified
VINE_GREEN = (0, 180, 136)

# Add green band title
def add_blue_band_title(slide, text, band_width=6.2):
    left, top, height = Inches(0.5), Inches(0.7), Inches(0.7)
    width = Inches(band_width)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(*VINE_GREEN)
    band.line.fill.background()
    title_box = slide.shapes.add_textbox(left + Inches(0.2), top, width-Inches(0.3), height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255,255,255)
    p.font.name = "Calibri"
    tf.word_wrap = True
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.margin_left = 0
    tf.margin_right = 0

def add_section_header(prs, title_text, subtitle_text="", color=VINE_GREEN):
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)
    slide.shapes.title.text = title_text
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
    slide.shapes.title.text_frame.paragraphs[0].font.name = 'Calibri'
    if subtitle_text:
        slide.placeholders[1].text = subtitle_text
        slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(22)
        slide.placeholders[1].text_frame.paragraphs[0].font.color.rgb = RGBColor(200,220,255)
    return slide

def add_title_slide(prs, title_text, subtitle_text, vine_logo_url):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(28, 34, 43)
    slide.shapes.title.text = title_text
    title = slide.shapes.title
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.name = 'Calibri Light'
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    slide.placeholders[1].text = subtitle_text
    ph = slide.placeholders[1]
    ph.text_frame.paragraphs[0].font.size = Pt(28)
    ph.text_frame.paragraphs[0].font.color.rgb = RGBColor(180, 200, 255)
    ph.text_frame.paragraphs[0].font.name = 'Calibri'
    img_data = requests.get(vine_logo_url).content
    img_stream = BytesIO(img_data)
    img_width = Inches(2.2)
    img_height = Inches(1.1)
    slide_width = prs.slide_width
    img_left = int((slide_width - img_width) / 2)
    img_top = Inches(0.2)
    slide.shapes.add_picture(img_stream, img_left, img_top, width=img_width, height=img_height)
    return slide

def add_custom_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_blue_band_title(slide, "Vine Overview", band_width=6.2)
    left = Inches(0.7); top = Inches(1.7); width = Inches(2.2); height = Inches(1)
    found_box = slide.shapes.add_textbox(left, top, width, height)
    tf = found_box.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "Founded in June 2012\n(Dom Hofmann, Rus Yusupov, Colin Kroll)\nAcquired by Twitter for ~$30M, Oct 2012"
    p1.font.size = Pt(16)
    p1.font.name = 'Calibri'
    p1.space_after = Pt(6)
    desc_box = slide.shapes.add_textbox(left, top + height + Inches(0.15), width, Inches(0.7))
    tfd = desc_box.text_frame
    pd = tfd.paragraphs[0]
    pd.text = "A short-form video app for sharing 6-second looping videos."
    pd.font.size = Pt(15)
    pd.font.name = 'Calibri'
    kb_width = Inches(4.5)
    kb_height = Inches(1.1)
    kb_left = Inches(3.1)
    kb_top = Inches(3.8)
    kf_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,kb_left,kb_top,kb_width,kb_height)
    kf_box.fill.solid()
    kf_box.fill.fore_color.rgb = RGBColor(*VINE_GREEN)
    kf_box.line.fill.background()
    kf_tf = kf_box.text_frame
    kf = kf_tf.paragraphs[0]
    kf.text = "Launch: Jan 2013 • Peak: ~200M users (2015) • Shutdown: Jan 2017"
    kf.font.size = Pt(17)
    kf.font.name = 'Calibri'
    kf.font.color.rgb = RGBColor(255,255,255)
    kf.alignment = PP_ALIGN.CENTER

def add_band_bullets_slide(prs, title, bullets, wider_title=False):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    band_width = 7.5 if wider_title else 6.2
    add_blue_band_title(slide, title, band_width=band_width)
    top = Inches(1.6)
    left = Inches(1.0)
    width = Inches(8.5)
    height = Inches(4.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        if isinstance(b, (list, tuple)):
            p.text = b[0]
            p.level = 0
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.name = 'Calibri'
            for child in b[1:]:
                c = tf.add_paragraph()
                c.text = child
                c.level = 1
                c.font.size = Pt(16)
                c.font.italic = True
                c.font.name = 'Calibri Light'
        else:
            p.text = b
            p.level = 0
            p.font.size = Pt(19)
            p.font.name = 'Calibri'
    return slide

def add_band_quote_slide(prs, quote, source, wider_title=False):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    band_width = 7.5 if wider_title else 6.2
    add_blue_band_title(slide, "What People Said", band_width=band_width)
    left = Inches(1.0)
    top = Inches(2.0)
    width = Inches(8.0)
    height = Inches(2.3)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 255)
    line = shape.line
    line.color.rgb = RGBColor(120, 170, 255)
    tf = shape.text_frame
    tf.text = f"“{quote}”"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.color.rgb = RGBColor(30,30,30)
    tf.paragraphs[0].font.name = 'Calibri Light'
    p2 = tf.add_paragraph()
    p2.text = f"— {source}"
    p2.font.size = Pt(13)
    p2.alignment = PP_ALIGN.RIGHT
    p2.font.name = 'Calibri'
    p2.font.color.rgb = RGBColor(80,80,80)
    return slide

def add_band_image_slide(prs, title_text, image_path, height_inches=4.0, wider_title=False):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    band_width = 7.5 if wider_title else 6.2
    add_blue_band_title(slide, title_text, band_width=band_width)
    left = Inches(1.0)
    top = Inches(1.7)
    slide.shapes.add_picture(image_path, left, top, height=Inches(height_inches))
    return slide

# Helper to add image to bottom/center of slide
def add_bottom_center_image(slide, prs, img_url, img_width_in=3.5, img_height_in=None, bottom_margin_in=0.5):
    img_data = requests.get(img_url).content
    img_stream = BytesIO(img_data)
    slide_width = prs.slide_width
    width = Inches(img_width_in)
    height = Inches(img_height_in) if img_height_in else None
    left = int((slide_width - width) / 2)
    shapes_heights = [s.top+s.height for s in slide.shapes]
    bottom = max(shapes_heights) if shapes_heights else 0
    top = prs.slide_height - (height if height else Inches(img_width_in/3)) - Inches(bottom_margin_in)
    slide.shapes.add_picture(img_stream, left, top, width=width, height=height)

# Chart data
months = ["Jan '13", "Jun '13", "Dec '13", "Dec '14", "Dec '15", "Jan '17"]
users_m = [1, 13, 40, 100, 200, 0]
engagement_pct = [85, 78, 72, 60, 35, 0]
def make_chart(x, y, title, yl, fname, color='#00b488'):
    fig = plt.figure(figsize=(6,3.5))
    plt.plot(x, y, marker='o', color=color, linewidth=3)
    plt.title(title, fontsize=16, weight='bold')
    plt.xlabel("Date", fontsize=12)
    plt.ylabel(yl, fontsize=12)
    plt.grid(True, which='both', linestyle='--', linewidth=0.6, alpha=0.7)
    for i, j in zip(x, y):
        plt.annotate(f"{j:.1f}", xy=(i, j), textcoords="offset points", xytext=(0, 7), ha='center', fontsize=11, color=color)
    plt.tight_layout()
    plt.savefig(fname, dpi=200)
    plt.close(fig)
chart1_path = "vine_users.png"
chart2_path = "vine_engagement.png"
make_chart(months, users_m, "Vine Registered Users (Millions) – Illustrative", "Users (M)", chart1_path)
make_chart(months, engagement_pct, "Vine Engagement Rate – Illustrative", "Engagement (%)", chart2_path, color='#00b488')

prs = Presentation()
add_title_slide(
    prs,
    "Why Vine Failed",
    "Auto-generated with Python for Microsoft PowerPynt",
    "https://variety.com/wp-content/uploads/2013/10/vine-logo.jpg?w=1000&h=563&crop=1"
)
add_section_header(prs, "Overview", "What is Vine?")
add_custom_overview_slide(prs)

slide = add_band_bullets_slide(
    prs, "Core Platform Features",
    [
        "6-second looping videos",
        "Easy uploading and browsing on mobile",
        "Instant sharing to Twitter/Facebook",
        "Simple interface, minimal editing features",
        "Allowed multiple 'takes' per vine"
    ]
)
add_bottom_center_image(
    slide, prs, 
    "https://techcrunch.com/wp-content/uploads/2013/07/capture-explore.jpg",
    img_width_in=3.2, img_height_in=None, bottom_margin_in=2
)

slide = add_band_bullets_slide(
    prs, "Key Milestones & Growth",
    [
        "Launch: Jan 2013 on iOS, later Android/Xbox",
        "Reached 40M users within a year",
        "Trendsetting among comedians, musicians, and meme creators",
        "Extremely young user base (teens, Gen Z)",
        "Acquired by Twitter before launch"
    ]
)
add_bottom_center_image(
    slide, prs, 
    "https://cdn.statcdn.com/Infographic/images/normal/1553.jpeg",
    img_width_in=4.2, img_height_in=None, bottom_margin_in=2
)

slide = add_band_bullets_slide(
    prs, "Peak Popularity",
    [
        "Attracted celebrities and brands for viral content",
        "Became a meme powerhouse, launching internet stars",
        "Many creators broke out to mainstream fame",
        "Community culture: 'Viners', collaborations, trends",
        "Daily active users peaked around 30M in 2014-2015"
    ],
    wider_title=True
)

slide = add_band_bullets_slide(
    prs, "Competitive Pressures",
    [
        "Instagram launched video ~6 months after Vine",
        "Snapchat took over ephemeral creativity",
        "Rise of Musical.ly/TikTok: longer content, more flexibility",
        "Other platforms had monetization built-in",
        "YouTube and IG attracted top Vine creators"
    ]
)

slide = add_band_bullets_slide(
    prs, "Business Model Issues",
    [
        "No revenue sharing or creator funds",
        "No ad network or influencer partnerships",
        "Extremely limited monetization options for top users",
        "Relied solely on parent (Twitter) for financial support",
        "No innovation in paid features or expansion"
    ],
    wider_title=True
)

slide = add_band_bullets_slide(
    prs, "Positioning",
    [
        "Short videos became less distinctive over time",
        "Lacked editing/effects: TikTok became more creative",
        "Brand did not adapt to creator needs",
        "Celebrity and advertising appeal faded in late years",
        "No response to algorithmic content discovery revolution"
    ],
    wider_title=True
)

add_band_image_slide(prs, "Growth Signal (Illustrative)", chart1_path, wider_title=True)
add_band_image_slide(prs, "Engagement Decline (Illustrative)", chart2_path, wider_title=True)

slide = add_band_bullets_slide(
    prs, "Execution & Org Gaps",
    [
        "Slow to adapt to feature requests (e.g., longer clips)",
        "Weak outreach/support for top creators",
        "Leadership turnover post-acquisition",
        "Poor integration with Twitter ecosystem",
        "Minimal feedback loops for user engagement"
    ],
    wider_title=True
)

slide = add_band_bullets_slide(
    prs, "Missed Opportunities",
    [
        "Could not pivot to trends like vines on YouTube",
        "Neglected tools for remixing and duets",
        "Few attempts at international growth"
    ]
)

slide = add_band_quote_slide(prs, "Vine made us stars, but didn’t help us make a living.", "Popular former Viner", wider_title=True)
add_bottom_center_image(
    slide, prs, 
    "https://www.360psg.com/content/images/Vine-vs-Instagram.jpg",
    img_width_in=3.3, img_height_in=None, bottom_margin_in=1
)

slide = add_band_bullets_slide(
    prs, "How Instagram and TikTok beat Vine",
    [
        "Longer videos and better editing tools",
        "Algorithmic content recommendation and discovery",
        "Direct and indirect monetization opportunities",
        "Richer social network, discovery, and sharing"
    ],
    wider_title=True
)

# Final slide: image at bottom, white font
slide = add_band_bullets_slide(
    prs, "Thank You",
    [
        "Generated with python-pptx + matplotlib",
        "Team: <your names here>",
        "GitHub repo: <link here>"
    ],
    wider_title=True
)
add_bottom_center_image(
    slide, prs,
    "https://media-cldnry.s-nbcnews.com/image/upload/t_nbcnews-fp-1200-630,f_auto,q_auto:best/rockcms/2022-01/220114-vine-eulogy-social-cs-dc5e98.png",
    img_width_in=4, img_height_in=None, bottom_margin_in=3
)
for shape in slide.shapes:
    if hasattr(shape, "text_frame"):
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)

prs.save("powerpynt_vine_finalfinal.pptx")
print("Saved presentation to: powerpynt_vine_finalfinal.pptx")